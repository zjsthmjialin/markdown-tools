"""Direct format-to-format converters that preserve original formatting."""

import os


class DirectConverter:
    """Handles direct conversion between formats without Markdown intermediate."""

    @staticmethod
    def pdf_to_docx(file_path: str, output_dir: str) -> dict:
        """Convert PDF to DOCX using PyMuPDF rendering + Tesseract OCR.
        Produces editable text without requiring Acrobat."""
        import fitz
        from docx import Document
        from docx.shared import Pt, Cm, Inches
        from ocr.tesseract_ocr import TesseractOCR

        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_path = os.path.join(output_dir, base_name + '.docx')
        os.makedirs(output_dir, exist_ok=True)

        ocr = TesseractOCR()
        doc = fitz.open(file_path)
        document = Document()

        # Remove default empty paragraph
        if document.paragraphs:
            p = document.paragraphs[0]._element
            p.getparent().remove(p)

        # Set A4 margins
        section = document.sections[0]
        section.page_width = Inches(8.27)
        section.page_height = Inches(11.69)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)

        total_pages = len(doc)
        for page_num in range(total_pages):
            page = doc[page_num]

            # Render at high DPI for OCR quality (300 DPI)
            mat = fitz.Matrix(300 / 72, 300 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")

            # OCR the page image
            text = ocr.extract_text_from_image(img_bytes)

            if text.strip():
                # Add as page title
                document.add_heading(f'第 {page_num + 1} 页', level=2)

                # Parse OCR text into paragraphs (split by newlines)
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue

                    # Try to detect and preserve table structure
                    # If line has multiple tabs or consistent spacing, treat as table row
                    if '\t' in line:
                        cells = line.split('\t')
                        table = document.add_table(rows=1, cols=len(cells))
                        table.style = 'Light Grid Accent 1'
                        row = table.rows[0]
                        for i, cell_text in enumerate(cells):
                            row.cells[i].text = cell_text.strip()
                    else:
                        p = document.add_paragraph(line)
                        p.paragraph_format.space_after = Pt(2)
                        p.paragraph_format.line_spacing = Pt(14)
            else:
                # No text found, add placeholder
                document.add_paragraph(f'[第 {page_num + 1} 页无文字内容]')

            # Page break (except last)
            if page_num < total_pages - 1:
                document.add_page_break()

        doc.close()
        document.save(output_path)
        return {'success': True, 'outputPath': output_path}

    @staticmethod
    def pdf_to_pptx(file_path: str, output_dir: str) -> dict:
        """Convert PDF to PPTX by rendering each page as a high-res image."""
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
        import io

        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_path = os.path.join(output_dir, base_name + '.pptx')
        os.makedirs(output_dir, exist_ok=True)

        doc = fitz.open(file_path)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for page_num in range(len(doc)):
            page = doc[page_num]
            mat = fitz.Matrix(300 / 72, 300 / 72)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")

            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0),
                                     Inches(13.333), Inches(7.5))

        doc.close()
        prs.save(output_path)
        return {'success': True, 'outputPath': output_path}

    @staticmethod
    def docx_to_pdf(file_path: str, output_dir: str) -> dict:
        """Convert DOCX to PDF using Microsoft Word COM."""
        import subprocess
        import tempfile

        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_path = os.path.join(output_dir, base_name + '.pdf')
        os.makedirs(output_dir, exist_ok=True)

        abs_input = os.path.abspath(file_path)
        abs_output = os.path.abspath(output_path)
        input_repr = repr(abs_input)
        output_repr = repr(abs_output)

        script = f'''
import os
import pythoncom
from win32com.client import Dispatch

input_path = {input_repr}
output_path = {output_repr}

pythoncom.CoInitialize()
try:
    word = Dispatch("Word.Application")
    word.Visible = False
    doc = word.Documents.Open(input_path, ConfirmConversions=False)
    doc.SaveAs(output_path, 17)
    doc.Close(False)
    word.Quit()
    print("OK")
finally:
    pythoncom.CoUninitialize()
'''

        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False, encoding='utf-8') as f:
            f.write(script)
            tmp_script = f.name

        try:
            result = subprocess.run(['python', tmp_script], capture_output=True,
                                    text=True, timeout=300)
            if result.returncode != 0 or 'OK' not in (result.stdout or ''):
                stderr = (result.stderr or '').encode('utf-8', errors='replace').decode('utf-8', errors='replace') if result.stderr else ''
                stdout = (result.stdout or '').encode('utf-8', errors='replace').decode('utf-8', errors='replace') if result.stdout else ''
                return {'success': False, 'error': f'DOCX转PDF失败: {stdout}{stderr}'[:500]}
            if not os.path.exists(output_path):
                return {'success': False, 'error': 'DOCX转PDF失败: 输出文件未生成'}
        finally:
            try:
                os.unlink(tmp_script)
            except Exception:
                pass

        return {'success': True, 'outputPath': output_path}

    @staticmethod
    def pptx_to_pdf(file_path: str, output_dir: str) -> dict:
        """Convert PPTX to PDF using PowerPoint COM."""
        import subprocess
        import tempfile

        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_path = os.path.join(output_dir, base_name + '.pdf')
        os.makedirs(output_dir, exist_ok=True)

        abs_input = os.path.abspath(file_path)
        abs_output = os.path.abspath(output_path)

        input_repr = repr(abs_input)
        output_repr = repr(abs_output)

        script = f'''
import os
import pythoncom
from win32com.client import Dispatch

input_path = {input_repr}
output_path = {output_repr}

pythoncom.CoInitialize()
try:
    powerpoint = Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
    presentation.SaveAs(output_path, 32)
    presentation.Close()
    powerpoint.Quit()
    print("OK")
finally:
    pythoncom.CoUninitialize()
'''

        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False, encoding='utf-8') as f:
            f.write(script)
            tmp_script = f.name

        try:
            result = subprocess.run(['python', tmp_script], capture_output=True,
                                    text=True, timeout=120)
            if result.returncode != 0 or 'OK' not in (result.stdout or ''):
                stderr = (result.stderr or '').encode('utf-8', errors='replace').decode('utf-8', errors='replace') if result.stderr else ''
                stdout = (result.stdout or '').encode('utf-8', errors='replace').decode('utf-8', errors='replace') if result.stdout else ''
                return {'success': False, 'error': f'PPTX转PDF失败: {stdout}{stderr}'[:500]}
            if not os.path.exists(output_path):
                return {'success': False, 'error': 'PPTX转PDF失败: 输出文件未生成'}
        finally:
            try:
                os.unlink(tmp_script)
            except Exception:
                pass

        return {'success': True, 'outputPath': output_path}

    _DIRECT_PAIRS = {
        ('.pdf', 'docx'): 'pdf_to_docx',
        ('.pdf', 'pptx'): 'pdf_to_pptx',
        ('.docx', 'pdf'): 'docx_to_pdf',
        ('.doc', 'pdf'): 'docx_to_pdf',
        ('.pptx', 'pdf'): 'pptx_to_pdf',
        ('.ppt', 'pdf'): 'pptx_to_pdf',
        ('.docx', 'docx'): None,
        ('.pdf', 'pdf'): None,
    }

    @staticmethod
    def can_direct_convert(source_ext: str, target_format: str) -> bool:
        return (source_ext, target_format) in DirectConverter._DIRECT_PAIRS

    @staticmethod
    def convert(file_path: str, source_ext: str, target_format: str, output_dir: str) -> dict:
        pair = (source_ext, target_format)
        if pair not in DirectConverter._DIRECT_PAIRS:
            return {'success': False, 'error': f'不支持的直接转换: {source_ext} -> {target_format}'}

        method_name = DirectConverter._DIRECT_PAIRS[pair]
        if method_name is None:
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            ext = '.' + target_format
            output_path = os.path.join(output_dir, base_name + ext)
            os.makedirs(output_dir, exist_ok=True)
            import shutil
            shutil.copy2(file_path, output_path)
            return {'success': True, 'outputPath': output_path}

        method = getattr(DirectConverter, method_name)
        return method(file_path, output_dir)