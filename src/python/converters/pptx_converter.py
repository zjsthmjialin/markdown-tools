from .base import BaseConverter, sanitize_text
from pptx import Presentation

class PptxConverter(BaseConverter):
    def get_supported_extensions(self) -> list[str]:
        return ['.pptx', '.ppt']

    def to_markdown(self, file_path: str, output_dir: str) -> str:
        prs = Presentation(file_path)
        markdown_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_parts.append(f"## 幻灯片 {slide_num}\n")

            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = sanitize_text(shape.text.strip())
                    if text:
                        markdown_parts.append(text)
                        markdown_parts.append('')

                if shape.has_table:
                    table_md = self._table_to_markdown(shape.table)
                    if table_md:
                        markdown_parts.append(table_md)
                        markdown_parts.append('')

        return '\n\n'.join(markdown_parts)

    def _table_to_markdown(self, table) -> str:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [sanitize_text(cell.text.strip()) for cell in row.cells]
            rows.append('| ' + ' | '.join(cells) + ' |')
            if i == 0:
                rows.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
        return '\n'.join(rows)