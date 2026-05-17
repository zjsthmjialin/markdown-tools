import re
import os
from xml.sax.saxutils import escape as xml_escape

# Use 'markdown' PyPI library via subprocess to avoid name collision
# with local 'markdown' package (src/python/markdown/).
# For direct in-process usage, we delegate to a helper.
from docx import Document
from docx.shared import Pt


def _render_markdown_to_html(md_text: str) -> str:
    """Render Markdown to HTML using the PyPI markdown library, avoiding
    the namespace collision with our local markdown package."""
    import sys
    import codecs

    # Save and remove ALL markdown-related modules from sys.modules
    _saved = {}
    for key in list(sys.modules.keys()):
        if key == 'markdown' or key.startswith('markdown.'):
            _saved[key] = sys.modules.pop(key)

    # Remove ALL sys.path entries that could resolve to the local markdown package.
    # Check by seeing if the path contains 'markdown' subdir with __init__.py
    # that is NOT in site-packages.
    _orig_path = sys.path[:]
    sys.path = [p for p in sys.path
                if not (os.path.isfile(os.path.join(p, 'markdown', '__init__.py'))
                        and 'site-packages' not in p)]

    try:
        import markdown as _md_lib
        result = _md_lib.markdown(md_text, extensions=['tables', 'fenced_code'])
        return result
    except UnicodeDecodeError:
        # Fallback: try to clean the input text
        cleaned = md_text.encode('utf-8', errors='replace').decode('utf-8', errors='replace')
        return _md_lib.markdown(cleaned, extensions=['tables', 'fenced_code'])
    finally:
        # Restore everything
        sys.path = _orig_path
        sys.modules.update(_saved)


def _is_table_separator(line: str) -> bool:
    """Check if a line is a markdown table separator like |---|---|"""
    return bool(re.match(r'^\|[-\s|]+\|$', line.strip()))


def _parse_table_row(line: str):
    """Parse a markdown table row into a list of cell values."""
    return [c.strip() for c in line.strip().split('|')[1:-1]]


class MarkdownToDocument:
    def to_pdf(self, markdown_content: str, output_path: str) -> None:
        """将 Markdown 转换为 PDF，使用 reportlab"""
        self._pdf_with_reportlab(markdown_content, output_path)

    def _pdf_with_reportlab(self, markdown_content: str, output_path: str) -> None:
        """使用 reportlab 生成 PDF（支持表格）"""
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        # 注册中文字体
        font_paths = [
            'C:/Windows/Fonts/msyh.ttc',    # 微软雅黑
            'C:/Windows/Fonts/simhei.ttf',   # 黑体
            'C:/Windows/Fonts/simsun.ttc',   # 宋体
        ]
        font_name = 'Helvetica'
        for fp in font_paths:
            if os.path.exists(fp):
                try:
                    pdfmetrics.registerFont(TTFont('ChineseFont', fp))
                    font_name = 'ChineseFont'
                    break
                except Exception:
                    pass

        doc = SimpleDocTemplate(output_path, pagesize=A4,
                                leftMargin=25*mm, rightMargin=25*mm,
                                topMargin=25*mm, bottomMargin=25*mm)
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(name='CNBody', fontName=font_name, fontSize=11, leading=18))
        styles.add(ParagraphStyle(name='CNH1', fontName=font_name, fontSize=20, leading=28, spaceAfter=12))
        styles.add(ParagraphStyle(name='CNH2', fontName=font_name, fontSize=16, leading=24, spaceAfter=10))
        styles.add(ParagraphStyle(name='CNH3', fontName=font_name, fontSize=13, leading=20, spaceAfter=8))

        elements = []
        table_buffer = []  # buffer for table rows

        def flush_table():
            """Flush buffered table rows into a reportlab Table element."""
            if not table_buffer:
                return
            # First row is header, rest are data
            num_cols = max(len(row) for row in table_buffer)
            # Pad rows to have equal columns
            for row in table_buffer:
                while len(row) < num_cols:
                    row.append('')
            table = Table(table_buffer)
            style_commands = [
                ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.91, 0.93, 0.96)),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                ('FONTNAME', (0, 0), (-1, 0), font_name),
                ('FONTSIZE', (0, 0), (-1, 0), 11),
                ('FONTNAME', (0, 1), (-1, -1), font_name),
                ('FONTSIZE', (0, 1), (-1, -1), 10),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.Color(0.96, 0.97, 0.98)]),
                ('TOPPADDING', (0, 0), (-1, -1), 4),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                ('LEFTPADDING', (0, 0), (-1, -1), 6),
                ('RIGHTPADDING', (0, 0), (-1, -1), 6),
            ]
            table.setStyle(TableStyle(style_commands))
            elements.append(table)
            elements.append(Spacer(1, 6))
            table_buffer.clear()

        for line in markdown_content.split('\n'):
            line = line.strip()
            if not line:
                flush_table()
                continue

            # 表格行
            if line.startswith('|'):
                # Skip separator rows
                if _is_table_separator(line):
                    continue
                table_buffer.append([xml_escape(c) for c in _parse_table_row(line)])
                continue
            else:
                flush_table()

            # 标题
            heading = re.match(r'^(#{1,6})\s+(.+)$', line)
            if heading:
                level = len(heading.group(1))
                text = xml_escape(heading.group(2))
                style = [styles['CNH3'], styles['CNH3'], styles['CNH3'],
                          styles['CNH2'], styles['CNH1'], styles['CNH1']][min(level, 6) - 1]
                elements.append(Paragraph(text, style))
                continue

            # 列表
            if re.match(r'^[-*]\s+', line):
                elements.append(Paragraph(f'• {xml_escape(line[2:])}', styles['CNBody']))
                continue

            # 普通段落
            elements.append(Paragraph(xml_escape(line), styles['CNBody']))

        # Flush any remaining table
        flush_table()

        if elements:
            doc.build(elements)

    def to_docx(self, markdown_content: str, output_path: str) -> None:
        # Clean control characters that break XML
        markdown_content = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', markdown_content)
        doc = Document()

        lines = markdown_content.split('\n')
        in_code_block = False
        code_lines = []
        table_buffer = []  # buffer for table rows

        def flush_table():
            """Flush buffered table rows into a Word table."""
            if not table_buffer:
                return
            num_cols = max(len(row) for row in table_buffer)
            # Pad rows to have equal columns
            for row in table_buffer:
                while len(row) < num_cols:
                    row.append('')
            table = doc.add_table(rows=len(table_buffer), cols=num_cols)
            table.style = 'Light Grid Accent 1'
            for i, row_data in enumerate(table_buffer):
                for j, cell_text in enumerate(row_data):
                    table.rows[i].cells[j].text = cell_text
            table_buffer.clear()

        for line in lines:
            # Code block handling
            if line.strip().startswith('```'):
                if in_code_block:
                    flush_table()
                    p = doc.add_paragraph()
                    run = p.add_run('\n'.join(code_lines))
                    run.font.name = 'Consolas'
                    run.font.size = Pt(10)
                    code_lines = []
                else:
                    flush_table()
                    in_code_block = True
                continue

            if in_code_block:
                code_lines.append(line)
                continue

            # Table row buffering
            if line.startswith('|'):
                if _is_table_separator(line):
                    continue
                table_buffer.append(_parse_table_row(line))
                continue
            else:
                flush_table()

            heading = re.match(r'^(#{1,6})\s+(.+)$', line)
            if heading:
                level = len(heading.group(1))
                text = heading.group(2)
                doc.add_heading(text, level=min(level, 9))
                continue

            if re.match(r'^[-*]\s+', line):
                doc.add_paragraph(line[2:], style='List Bullet')
                continue

            if line.strip():
                doc.add_paragraph(line)

        # Flush any remaining table
        flush_table()

        doc.save(output_path)

    def to_pptx(self, markdown_content: str, output_path: str) -> None:
        """Convert Markdown to PPTX. Each H1 starts a new slide."""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Split content by H1 headings into slides
        sections = re.split(r'\n(?=^#\s)', markdown_content, flags=re.MULTILINE)

        for section in sections:
            lines = section.strip().split('\n')
            if not lines:
                continue

            slide_layout = prs.slide_layouts[6]  # blank layout
            slide = prs.slides.add_slide(slide_layout)

            left = Inches(1)
            top = Inches(0.5)
            width = Inches(11.333)
            height = Inches(6.5)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True

            first = True
            table_buffer = []  # buffer for table rows

            def flush_table():
                """Flush buffered table rows into a PowerPoint table."""
                if not table_buffer:
                    return
                num_rows = len(table_buffer)
                num_cols = max(len(row) for row in table_buffer)
                # Pad rows to have equal columns
                for row in table_buffer:
                    while len(row) < num_cols:
                        row.append('')

                tbl_width = Inches(11.333)
                row_height = Inches(0.4)
                tbl_height = row_height * num_rows
                tbl_left = Inches(1)
                tbl_top = Inches(0.5)

                table_shape = slide.shapes.add_table(
                    num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height
                )
                table = table_shape.table

                for i, row_data in enumerate(table_buffer):
                    for j, cell_text in enumerate(row_data):
                        cell = table.cell(i, j)
                        cell.text = cell_text
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(14)

                table_buffer.clear()

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Table row buffering
                if line.startswith('|'):
                    if _is_table_separator(line):
                        continue
                    table_buffer.append(_parse_table_row(line))
                    continue
                else:
                    flush_table()

                h_match = re.match(r'^(#{1,6})\s+(.+)$', line)
                if h_match:
                    level = len(h_match.group(1))
                    text = h_match.group(2)
                    p = tf.add_paragraph() if not first else tf.paragraphs[0]
                    first = False
                    p.text = text
                    if level == 1:
                        p.font.size = Pt(36)
                        p.font.bold = True
                    elif level == 2:
                        p.font.size = Pt(28)
                        p.font.bold = True
                    else:
                        p.font.size = Pt(22)
                        p.font.bold = True
                    p.space_after = Pt(12)
                    continue

                if re.match(r'^[-*]\s+', line):
                    p = tf.add_paragraph() if not first else tf.paragraphs[0]
                    first = False
                    p.text = f'  {line[2:]}'
                    p.font.size = Pt(18)
                    p.space_after = Pt(6)
                    continue

                # Regular paragraph
                p = tf.add_paragraph() if not first else tf.paragraphs[0]
                first = False
                p.text = line
                p.font.size = Pt(18)
                p.space_after = Pt(8)

            # Flush any remaining table
            flush_table()

        prs.save(output_path)

    def to_html(self, markdown_content: str, output_path: str) -> None:
        html_content = self._markdown_to_html(markdown_content)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)

    def _markdown_to_html(self, markdown_content: str) -> str:
        return f'''<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
body {{ font-family: -apple-system, "Noto Sans SC", "Microsoft YaHei", sans-serif;
        max-width: 800px; margin: 0 auto; padding: 20px; line-height: 1.8; }}
pre {{ background: #f5f5f5; padding: 16px; border-radius: 8px; overflow-x: auto; }}
code {{ font-family: "JetBrains Mono", "Consolas", monospace; }}
table {{ border-collapse: collapse; width: 100%; }}
th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
th {{ background: #f5f5f5; }}
img {{ max-width: 100%; }}
</style>
</head>
<body>
{_render_markdown_to_html(markdown_content)}
</body>
</html>'''
